#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import sys
import json
import os
import base64
import tempfile
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import logging
logging.basicConfig(
    level=logging.INFO,
    format='%(levelname)s: %(message)s'
)
logger = logging.getLogger('basket_order_generator')

# Set constants for A4 landscape dimensions (in inches)
A4_WIDTH = 11.69  # 297mm
A4_HEIGHT = 8.27  # 210mm

def create_basket_order_slide(prs, order_data):
    logger.info(f"Creating slide for order: {order_data}")
    """Create a slide for a basket order"""
    # Add a blank slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide dimensions to A4 landscape
    prs.slide_width = Inches(A4_WIDTH)
    prs.slide_height = Inches(A4_HEIGHT)
    
    # Calculate margins and positions
    margin = Inches(0.5)
    
    # 添加白色背景
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # ----- TOP SECTION -----
    # Top row items: Date, Order Number, Color, Icon, Position
    
    # Date (生成日期)
    date_box = slide.shapes.add_textbox(margin, margin, Inches(3), Inches(0.5))
    date_text = date_box.text_frame
    date_p = date_text.paragraphs[0]
    date_p.text = order_data.get('date', '')
    date_p.font.bold = True
    date_p.font.size = Pt(24)
    date_p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
    logger.info(f"Added date: {date_p.text}")
    
    # Order Number (订单号)
    order_box_x = margin + Inches(3.5)
    order_box = slide.shapes.add_textbox(order_box_x, margin, Inches(4), Inches(0.5))
    order_text = order_box.text_frame
    order_p = order_text.paragraphs[0]
    order_p.text = str(order_data.get('orderNumber', ''))  # Convert to string
    order_p.font.bold = True
    order_p.font.size = Pt(20)
    order_p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
    logger.info(f"Added order number: {order_p.text}")
    
    # Color (毛线颜色)
    color_text = "默认颜色"
    if order_data.get('color'):
        color_text = order_data.get('color')
    
    color_box = slide.shapes.add_textbox(order_box_x + Inches(4.5), margin, Inches(2), Inches(0.5))
    color_text_frame = color_box.text_frame
    color_p = color_text_frame.paragraphs[0]
    color_p.text = color_text
    color_p.font.bold = True
    color_p.font.size = Pt(20)
    color_p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
    logger.info(f"Added color: {color_p.text}")
    
    # SKU (顶部展示)
    sku_text = f"SKU: {order_data.get('sku', '')}"
    sku_top_box = slide.shapes.add_textbox(margin, margin + Inches(0.6), Inches(5), Inches(0.5))
    sku_top_text = sku_top_box.text_frame
    sku_top_p = sku_top_text.paragraphs[0]
    sku_top_p.text = sku_text
    sku_top_p.font.bold = True
    sku_top_p.font.size = Pt(20)
    sku_top_p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
    logger.info(f"Added SKU at top: {sku_top_p.text}")
    
    # Icon (图标)
    icon_box = slide.shapes.add_textbox(margin + Inches(5.5), margin + Inches(0.6), Inches(3), Inches(0.5))
    icon_text = icon_box.text_frame
    icon_p = icon_text.paragraphs[0]
    icon_p.text = order_data.get('icon', '')
    icon_p.font.bold = True
    icon_p.font.size = Pt(20)
    icon_p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
    logger.info(f"Added icon: {icon_p.text}")
    
    # Position (一单多买的序号)
    position_box = slide.shapes.add_textbox(prs.slide_width - margin - Inches(1.5), margin, Inches(1.5), Inches(0.5))
    position_text = position_box.text_frame
    position_p = position_text.paragraphs[0]
    position_p.text = order_data.get('position', '')
    position_p.alignment = PP_ALIGN.RIGHT
    position_p.font.bold = True
    position_p.font.size = Pt(20)
    position_p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
    logger.info(f"Added position: {position_p.text}")
    
    # ----- MIDDLE SECTION -----
    # Recipient Name (收件人)
    recipient_box = slide.shapes.add_textbox(margin, margin + Inches(1.5), prs.slide_width - margin * 2, Inches(1.5))
    recipient_text = recipient_box.text_frame
    recipient_p = recipient_text.paragraphs[0]
    recipient_p.text = order_data.get('recipientName', '')
    recipient_p.alignment = PP_ALIGN.CENTER
    recipient_p.font.bold = True
    recipient_p.font.size = Pt(36)
    recipient_p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
    logger.info(f"Added recipient name: {recipient_p.text}")
    
    # ----- MAIN CONTENT -----
    # Custom Value (定制内容) - Large text that takes up most of the slide
    value_box = slide.shapes.add_textbox(margin, margin + Inches(3), prs.slide_width - margin * 2, Inches(3.5))
    value_text = value_box.text_frame
    value_text.word_wrap = True
    value_p = value_text.paragraphs[0]
    value_p.text = order_data.get('customName', '')
    value_p.alignment = PP_ALIGN.CENTER
    value_p.font.bold = True
    value_p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
    
    # Adjust font size based on text length to ensure it fits well
    text_len = len(value_p.text)
    if text_len > 30:
        value_p.font.size = Pt(36)
    elif text_len > 15:
        value_p.font.size = Pt(48)
    else:
        value_p.font.size = Pt(72)
    logger.info(f"Added custom name: {value_p.text} with font size: {value_p.font.size.pt}")
    
    # ----- BOTTOM SECTION -----
    # Shop name and SKU info at the bottom
    bottom_box = slide.shapes.add_textbox(margin, prs.slide_height - margin - Inches(0.5), prs.slide_width - margin * 2, Inches(0.5))
    bottom_text_frame = bottom_box.text_frame
    bottom_p = bottom_text_frame.paragraphs[0]
    
    # Combine shop name and SKU
    shop_name = order_data.get('shopName', '')
    
    if shop_name:
        bottom_text = shop_name
    else:
        bottom_text = ""
    
    # Add quantity if more than 1
    if order_data.get('quantity') and order_data.get('quantity') > 1:
        if bottom_text:
            bottom_text += f" | 数量: {order_data.get('quantity')}"
        else:
            bottom_text = f"数量: {order_data.get('quantity')}"
    
    bottom_p.text = bottom_text
    bottom_p.alignment = PP_ALIGN.CENTER
    bottom_p.font.size = Pt(16)
    bottom_p.font.bold = True
    bottom_p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
    logger.info(f"Added bottom text: {bottom_p.text}")
    
    # Return the slide
    return slide

def process_json_data(orders_data):
    """Process JSON data and generate PPT slides"""
    try:
        # Create a new presentation
        prs = Presentation()
        
        # Create slides for each order in the data
        slides_created = 0
        errors = []
        
        for index, order_data in enumerate(orders_data):
            try:
                # Create slide for this order
                create_basket_order_slide(prs, order_data)
                slides_created += 1
            except Exception as e:
                errors.append(f"Error processing order {index + 1}: {str(e)}")
        
        # Save the presentation to a temporary file
        output_path = os.path.join(tempfile.gettempdir(), 'basket_orders.pptx')
        prs.save(output_path)
        
        # Read the file and encode it as base64
        with open(output_path, 'rb') as file:
            ppt_data = base64.b64encode(file.read()).decode('utf-8')
        
        # Clean up temporary file
        os.unlink(output_path)
        
        # Return success result with the PPT data
        result = {
            'success': True,
            'message': f'Successfully generated {slides_created} slides',
            'slides_created': slides_created,
            'data': ppt_data,
            'errors': errors if errors else None
        }
        
    except Exception as e:
        # Return error result
        result = {
            'success': False,
            'message': f'Error generating PPT: {str(e)}',
            'error': traceback.format_exc()
        }
    
    return result

def main():
    """Main function to handle command line input and output"""
    try:
        # Read JSON input from stdin
        input_data = json.loads(sys.stdin.read())
        
        # Get the output path if provided
        output_path = input_data.get('outputPath')
        
        # Check for the data (could be in 'excelData' for backward compatibility or 'orderData' for new format)
        json_data = input_data.get('orderData') or input_data.get('excelData')
        
        if not json_data:
            result = {
                'success': False,
                'message': 'No data provided'
            }
        else:
            # Decode base64 data if provided
            try:
                decoded_data = base64.b64decode(json_data).decode('utf-8')
                orders_data = json.loads(decoded_data)
                
                # Process the orders data to generate PPT
                result = process_json_data(orders_data)
                
                # If output path is provided and PPT data is available, write directly to file
                if output_path and result.get('success') and result.get('data'):
                    try:
                        ppt_dir = os.path.dirname(output_path)
                        if not os.path.exists(ppt_dir):
                            os.makedirs(ppt_dir)
                            
                        with open(output_path, 'wb') as f:
                            f.write(base64.b64decode(result['data']))
                        result['message'] += f" PPT saved to {output_path}"
                    except Exception as e:
                        result['message'] += f" (Warning: Failed to save to {output_path}: {str(e)})"
                
            except Exception as e:
                result = {
                    'success': False,
                    'message': f'Error decoding or parsing data: {str(e)}',
                    'error': traceback.format_exc()
                }
        
        # Output the result as JSON
        print(json.dumps(result))
        
    except Exception as e:
        # Return error result if there's an exception
        result = {
            'success': False,
            'message': f'Error: {str(e)}',
            'error': traceback.format_exc()
        }
        print(json.dumps(result))

if __name__ == '__main__':
    main() 